import hashlib
import hmac
import io
import logging
import uuid
from typing import List, Optional

from fastapi import (
    APIRouter,
    Depends,
    File,
    Form,
    HTTPException,
    Request,
    UploadFile,
    status,
)
from fastapi.responses import JSONResponse, Response, StreamingResponse
from sqlalchemy import or_
from sqlalchemy.orm import Session

from app.api.deps import get_current_user
from app.core.audit import audit_admin
from app.core.config import settings
from app.core.database import get_db
from app.core.ranges import parse_range
from app.models.models import Note, NoteFileType, NotesPurchase, User, UserRole
from app.schemas.schemas import (
    AdminGrantAccessRequest,
    NoteAdminResponse,
    NoteResponse,
    NotesAccessStatus,
    NotesPurchaseGrantResponse,
    NotesPurchaseOrderResponse,
    NotesPurchaseVerifyRequest,
)
import requests

from app.services.study_material_service import StudyMaterialService
from app.services.google_drive_service import GoogleDriveError

logger = logging.getLogger(__name__)

router = APIRouter()

# ---------------------------------------------------------------------------
# constants
# ---------------------------------------------------------------------------

ALLOWED_EXTENSIONS = {e.value for e in NoteFileType}

CONTENT_TYPE_MAP = {
    "pdf": "application/pdf",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "png": "image/png",
    "txt": "text/plain; charset=utf-8",
    "doc": "application/msword",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "ppt": "application/vnd.ms-powerpoint",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------


def _has_notes_access(user: User, db: Session) -> bool:
    if user.role == UserRole.ADMIN:
        return True
    if user.has_notes_access:
        return True
    purchase = (
        db.query(NotesPurchase)
        .filter(NotesPurchase.user_id == user.id, NotesPurchase.is_active == True)
        .first()
    )
    return purchase is not None


def _require_notes_access(user: User, db: Session):
    if not _has_notes_access(user, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Purchase Notes Pack to access this content",
        )


def _require_admin(user: User):
    if user.role != UserRole.ADMIN:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Admin access required")


def _get_file_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _spooled_size(file: UploadFile) -> int:
    """Size of an UploadFile without reading it into memory. Starlette spools
    uploads to a temp file past 1 MB, so we can measure via seek/tell and then
    stream the file object straight to Drive (bounded memory)."""
    f = file.file
    f.seek(0, 2)  # SEEK_END
    size = f.tell()
    f.seek(0)
    return size


def _extract_docx_content(file_bytes: bytes) -> dict:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                style = para.style.name if para.style else ""
                paragraphs.append({"heading": style.startswith("Heading"), "text": text})
        return {"type": "extracted_text", "paragraphs": paragraphs}
    except Exception as e:
        return {"type": "extracted_text", "paragraphs": [{"heading": False, "text": f"[Could not extract content: {e}]"}]}


def _extract_pptx_content(file_bytes: bytes) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            title_text = ""
            content_parts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape.shape_type == 13 or (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format
                    and shape.placeholder_format.idx == 0
                ):
                    title_text = text
                else:
                    content_parts.append(text)
            slides.append({"slide_num": i, "title": title_text, "content": "\n".join(content_parts)})
        return {"type": "extracted_slides", "slides": slides}
    except Exception as e:
        return {"type": "extracted_slides", "slides": [{"slide_num": 1, "title": "Error", "content": f"[Could not extract content: {e}]"}]}


# ---------------------------------------------------------------------------
# ACCESS STATUS
# ---------------------------------------------------------------------------


@router.get("/access-status", response_model=NotesAccessStatus)
def get_access_status(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    is_admin = current_user.role == UserRole.ADMIN
    has_access = _has_notes_access(current_user, db)
    purchased_at = None
    if not is_admin and has_access and not current_user.has_notes_access:
        purchase = (
            db.query(NotesPurchase)
            .filter(NotesPurchase.user_id == current_user.id, NotesPurchase.is_active == True)
            .first()
        )
        if purchase:
            purchased_at = purchase.purchased_at
    return NotesAccessStatus(has_access=has_access, is_admin=is_admin, purchased_at=purchased_at)


# ---------------------------------------------------------------------------
# PURCHASE
# ---------------------------------------------------------------------------


@router.post("/purchase")
def initiate_purchase(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    # Already has access
    if _has_notes_access(current_user, db):
        return {"status": "already_granted", "dev_mode": settings.DEVELOPMENT_MODE}

    if settings.DEVELOPMENT_MODE and not settings.is_production:
        # Auto-grant in dev
        purchase = NotesPurchase(
            user_id=current_user.id,
            razorpay_order_id="DEV_MODE",
            # Per-user sentinel (not a bare "DEV_MODE") so the production
            # unique index on razorpay_payment_id never collides in dev.
            razorpay_payment_id=f"DEV_{current_user.id}",
            amount=settings.NOTES_PRICE / 100,
            is_active=True,
        )
        db.add(purchase)
        db.commit()
        db.refresh(purchase)
        print(f"[DEV MODE] Notes access auto-granted to user {current_user.email}")
        return NotesPurchaseGrantResponse(status="granted", dev_mode=True, purchased_at=purchase.purchased_at)

    # Production: create Razorpay order
    if not settings.RAZORPAY_KEY_ID or not settings.RAZORPAY_KEY_SECRET:
        raise HTTPException(status_code=500, detail="Payment gateway not configured")

    url = "https://api.razorpay.com/v1/orders"
    auth = (settings.RAZORPAY_KEY_ID, settings.RAZORPAY_KEY_SECRET)
    payload = {
        "amount": settings.NOTES_PRICE,
        "currency": "INR",
        "receipt": f"notes_{str(current_user.id)[:8]}",
        "notes": {"user_id": str(current_user.id), "product": "notes_pack"},
    }
    try:
        resp = requests.post(url, json=payload, auth=auth, timeout=15)
    except requests.RequestException:
        raise HTTPException(status_code=503, detail="Payment gateway unreachable. Please try again.")
    if resp.status_code != 200:
        raise HTTPException(status_code=400, detail="Failed to create payment order")
    try:
        data = resp.json()
    except ValueError:
        raise HTTPException(status_code=502, detail="Invalid response from payment gateway")
    return NotesPurchaseOrderResponse(
        order_id=data["id"],
        amount=data["amount"],
        currency=data["currency"],
        key_id=settings.RAZORPAY_KEY_ID,
        dev_mode=False,
    )


@router.post("/purchase/verify", response_model=NotesPurchaseGrantResponse)
def verify_purchase(
    verify_in: NotesPurchaseVerifyRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if _has_notes_access(current_user, db):
        return NotesPurchaseGrantResponse(status="already_granted")

    if not (settings.DEVELOPMENT_MODE and not settings.is_production):
        generated = hmac.new(
            settings.RAZORPAY_KEY_SECRET.encode(),
            f"{verify_in.razorpay_order_id}|{verify_in.razorpay_payment_id}".encode(),
            hashlib.sha256,
        ).hexdigest()
        if not hmac.compare_digest(generated, verify_in.razorpay_signature):
            raise HTTPException(status_code=400, detail="Payment signature verification failed")

    # Bind each Razorpay payment to exactly one account. Without this a valid
    # (order|payment|signature) triple could be replayed to unlock OTHER users
    # (a single payment granting unlimited accounts), and re-submitting would
    # create duplicate purchase rows.
    existing = (
        db.query(NotesPurchase)
        .filter(NotesPurchase.razorpay_payment_id == verify_in.razorpay_payment_id)
        .first()
    )
    if existing:
        if existing.user_id == current_user.id:
            return NotesPurchaseGrantResponse(
                status="already_granted", dev_mode=settings.DEVELOPMENT_MODE,
                purchased_at=existing.purchased_at,
            )
        raise HTTPException(status_code=400, detail="This payment has already been used.")

    purchase = NotesPurchase(
        user_id=current_user.id,
        razorpay_order_id=verify_in.razorpay_order_id,
        razorpay_payment_id=verify_in.razorpay_payment_id,
        amount=settings.NOTES_PRICE / 100,
        is_active=True,
    )
    db.add(purchase)
    db.commit()
    db.refresh(purchase)
    return NotesPurchaseGrantResponse(status="granted", dev_mode=False, purchased_at=purchase.purchased_at)


# ---------------------------------------------------------------------------
# USER — LIST & METADATA
# ---------------------------------------------------------------------------


@router.get("", response_model=List[NoteResponse])
def list_notes(
    category: Optional[str] = None,
    search: Optional[str] = None,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """Return active notes metadata. Access check is NOT required here — users browse freely."""
    query = db.query(Note).filter(Note.is_active == True)
    if category:
        query = query.filter(Note.category == category)
    if search:
        term = f"%{search}%"
        query = query.filter(
            or_(Note.title.ilike(term), Note.description.ilike(term), Note.tags.ilike(term))
        )
    return query.order_by(Note.uploaded_at.desc()).all()


@router.get("/categories")
def list_categories(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    rows = db.query(Note.category).filter(Note.is_active == True).distinct().all()
    return [r[0] for r in rows]


@router.get("/{note_id}", response_model=NoteResponse)
def get_note(
    note_id: uuid.UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    note = db.query(Note).filter(Note.id == note_id, Note.is_active == True).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    return note


# ---------------------------------------------------------------------------
# VIEWER — SECURE FILE STREAMING
# ---------------------------------------------------------------------------


@router.get("/{note_id}/viewer")
def view_note(
    note_id: uuid.UUID,
    request: Request,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _require_notes_access(current_user, db)

    note = db.query(Note).filter(Note.id == note_id, Note.is_active == True).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    if not note.google_drive_file_id:
        raise HTTPException(status_code=404, detail="File is not available.")

    ft = note.file_type.value if isinstance(note.file_type, NoteFileType) else note.file_type
    svc = StudyMaterialService()

    # Office formats — extract server-side and return JSON (small payloads).
    if ft in ("docx", "doc", "pptx", "ppt"):
        try:
            file_bytes = svc.read(note.google_drive_file_id)
        except GoogleDriveError as e:
            raise HTTPException(status_code=e.status_code, detail=e.message)
        content = (
            _extract_docx_content(file_bytes)
            if ft in ("docx", "doc")
            else _extract_pptx_content(file_bytes)
        )
        return JSONResponse(content=content, headers={"Cache-Control": "no-store"})

    content_type = CONTENT_TYPE_MAP.get(ft, "application/octet-stream")
    base_headers = {
        "Content-Disposition": "inline",
        "Accept-Ranges": "bytes",
        "Cache-Control": "no-store, no-cache, must-revalidate",
        "Pragma": "no-cache",
        "X-Content-Type-Options": "nosniff",
    }

    # Range request → serve only the requested bytes straight from Drive, so the
    # PDF viewer loads pages progressively instead of pulling the whole file.
    range_header = request.headers.get("range")
    try:
        if range_header:
            # Use the size already stored in MSSQL — avoids a Drive metadata
            # round-trip on every range request (pdfrx issues many per document).
            total = note.file_size or svc.drive.get_file_size(note.google_drive_file_id)
            rng = parse_range(range_header, total)
            if rng is None:
                return Response(
                    status_code=416,
                    headers={"Content-Range": f"bytes */{total}", "Accept-Ranges": "bytes"},
                )
            start, end = rng
            chunk = svc.drive.download_range(note.google_drive_file_id, start, end)
            headers = dict(base_headers)
            headers["Content-Range"] = f"bytes {start}-{end}/{total}"
            headers["Content-Length"] = str(len(chunk))
            return Response(
                content=chunk, status_code=206, media_type=content_type, headers=headers
            )

        # No range → stream the whole file in chunks (bounded server memory).
        total = note.file_size or svc.drive.get_file_size(note.google_drive_file_id)
    except GoogleDriveError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)

    headers = dict(base_headers)
    if total:
        headers["Content-Length"] = str(total)
    return StreamingResponse(
        svc.drive.stream_file(note.google_drive_file_id),
        media_type=content_type,
        headers=headers,
    )


# ---------------------------------------------------------------------------
# ADMIN — UPLOAD
# ---------------------------------------------------------------------------


@router.post("/upload", response_model=NoteAdminResponse)
async def upload_note(
    title: str = Form(...),
    category: str = Form(...),
    description: Optional[str] = Form(None),
    tags: Optional[str] = Form(None),
    is_free: bool = Form(False),
    price: float = Form(0),
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _require_admin(current_user)

    ext = _get_file_extension(file.filename or "")
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"File type .{ext} not allowed. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}",
        )

    file_size = _spooled_size(file)
    max_bytes = settings.STUDY_MATERIAL_MAX_FILE_SIZE_MB * 1024 * 1024
    if file_size > max_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Maximum {settings.STUDY_MATERIAL_MAX_FILE_SIZE_MB} MB.",
        )

    content_type = CONTENT_TYPE_MAP.get(ext, "application/octet-stream")
    try:
        # Stream the disk-backed upload straight to Drive (bounded memory).
        file_id, folder_id = StudyMaterialService().store_stream(
            category=category,
            file_name=file.filename or f"note.{ext}",
            mime_type=content_type,
            fileobj=file.file,
        )
    except GoogleDriveError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)

    db_note = Note(
        title=title.strip(),
        description=description.strip() if description else None,
        category=category.strip(),
        tags=tags.strip() if tags else None,
        file_name=file.filename,
        google_drive_file_id=file_id,
        google_drive_folder_id=folder_id,
        mime_type=content_type,
        file_type=NoteFileType(ext),
        file_size=file_size,
        is_free=is_free,
        price=price,
        uploaded_by=current_user.id,
        is_active=True,
    )
    try:
        db.add(db_note)
        db.commit()
        db.refresh(db_note)
    except Exception:
        # Never orphan a file in Drive if the metadata insert fails.
        try:
            StudyMaterialService().remove(file_id)
        except Exception:
            pass
        raise
    return db_note


# ---------------------------------------------------------------------------
# ADMIN — EDIT METADATA
# ---------------------------------------------------------------------------


@router.put("/{note_id}", response_model=NoteAdminResponse)
async def update_note(
    note_id: uuid.UUID,
    title: Optional[str] = Form(None),
    category: Optional[str] = Form(None),
    description: Optional[str] = Form(None),
    tags: Optional[str] = Form(None),
    is_free: Optional[bool] = Form(None),
    price: Optional[float] = Form(None),
    file: Optional[UploadFile] = File(None),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _require_admin(current_user)

    note = db.query(Note).filter(Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")

    if title is not None:
        note.title = title.strip()
    if category is not None:
        note.category = category.strip()
    if description is not None:
        note.description = description.strip() or None
    if tags is not None:
        note.tags = tags.strip() or None
    if is_free is not None:
        note.is_free = is_free
    if price is not None:
        note.price = price

    if file is not None and file.filename:
        ext = _get_file_extension(file.filename)
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"File type .{ext} not allowed")
        file_size = _spooled_size(file)
        max_bytes = settings.STUDY_MATERIAL_MAX_FILE_SIZE_MB * 1024 * 1024
        if file_size > max_bytes:
            raise HTTPException(
                status_code=413,
                detail=f"File too large. Maximum {settings.STUDY_MATERIAL_MAX_FILE_SIZE_MB} MB.",
            )
        content_type = CONTENT_TYPE_MAP.get(ext, "application/octet-stream")
        svc = StudyMaterialService()
        try:
            new_file_id, new_folder_id = svc.store_stream(
                category=note.category,
                file_name=file.filename,
                mime_type=content_type,
                fileobj=file.file,
            )
        except GoogleDriveError as e:
            raise HTTPException(status_code=e.status_code, detail=e.message)
        # Best-effort removal of the previous Drive file after the new one lands.
        if note.google_drive_file_id:
            try:
                svc.remove(note.google_drive_file_id)
            except Exception:
                pass
        note.file_name = file.filename
        note.google_drive_file_id = new_file_id
        note.google_drive_folder_id = new_folder_id
        note.mime_type = content_type
        note.file_type = NoteFileType(ext)
        note.file_size = file_size

    db.commit()
    db.refresh(note)
    return note


# ---------------------------------------------------------------------------
# ADMIN — DELETE
# ---------------------------------------------------------------------------


@router.delete("/{note_id}")
def delete_note(
    note_id: uuid.UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _require_admin(current_user)

    note = db.query(Note).filter(Note.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")

    # Remove the Drive file first (404 is treated as already gone). A non-404
    # Drive failure is logged but does not block DB removal: a stray Drive file
    # is a minor storage leak, whereas a lingering row keeps a "deleted" note
    # visible in the admin list.
    if note.google_drive_file_id:
        try:
            StudyMaterialService().remove(note.google_drive_file_id)
        except GoogleDriveError as e:
            logger.warning("Drive delete failed for note %s: %s", note_id, e.message)

    # Hard delete — no soft-delete rows left behind.
    db.delete(note)
    db.commit()
    audit_admin("note.delete", actor=current_user, target_id=note_id)
    return {"status": "deleted"}


# ---------------------------------------------------------------------------
# ADMIN — LIST ALL NOTES
# ---------------------------------------------------------------------------


@router.get("/admin/all", response_model=List[NoteAdminResponse])
def admin_list_notes(
    category: Optional[str] = None,
    search: Optional[str] = None,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _require_admin(current_user)

    # Never return deleted notes (defence-in-depth; hard delete already removes rows).
    query = db.query(Note).filter(Note.is_active == True)
    if category:
        query = query.filter(Note.category == category)
    if search:
        term = f"%{search}%"
        query = query.filter(
            or_(Note.title.ilike(term), Note.description.ilike(term), Note.tags.ilike(term))
        )
    return query.order_by(Note.uploaded_at.desc()).all()


# ---------------------------------------------------------------------------
# ADMIN — GRANT / REVOKE ACCESS
# ---------------------------------------------------------------------------


@router.post("/admin/grant-access")
def admin_grant_access(
    req: AdminGrantAccessRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    _require_admin(current_user)

    user = db.query(User).filter(User.id == req.user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")

    user.has_notes_access = req.grant
    db.commit()
    audit_admin(
        "notes.grant_access" if req.grant else "notes.revoke_access",
        actor=current_user, target_id=req.user_id,
    )
    return {"user_id": str(req.user_id), "has_notes_access": req.grant}
